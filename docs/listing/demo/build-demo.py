"""
Build the deck the Marketplace screenshots are taken from.

    pip install python-pptx && python3 docs/listing/demo/build-demo.py

A SEPARATE deck from `test-kit/`, and the separation is the point. The kit's
template exists to break the engine: its photograph is a blue rectangle with
four yellow dots positioned to prove a cover-crop keeps the right axis, its
chart values are 18 and 42 because two numbers are enough to tell a swap from a
fill, and its first slide is a wall of instructions to whoever runs the round.
Every one of those choices is right for a test and wrong for a store listing,
where Microsoft asks for images that "show real content rather than an empty
document" and a prospective customer reads the crop fixture as a broken image.

The kit's deck also cannot simply be tidied up: `test/test-kit.test.ts` merges
the committed file and checks thirteen things about the result, so editing it to
look nicer breaks the round it exists to run.

Run rarely, by hand. Not wired into npm or CI, for the same reason
`test-kit/build-template.py` is not: it needs Python and a library this project
does not otherwise use, and the deck it writes is committed beside it.

WHAT IS DELIBERATELY NOT HERE: a picture placeholder. Merging pictures is a real
feature and it is in the listing description, but a screenshot of it needs a
photograph that looks like a photograph. Anything this script can draw would be
a coloured rectangle, which is the exact failure the kit's deck already has. A
shot of text, numbers, dates and a chart merging is honest; a shot with a fake
photograph in it is the placeholder problem again wearing a different colour.
"""

import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "Quarterly-business-review.pptx")

# Neutral and not the product's own palette: the pane is navy and orange, and a
# deck in the same two colours makes the screenshot look like one thing rather
# than an add-in working on somebody's file.
INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6B, 0x72, 0x80)
TEAL = RGBColor(0x0F, 0x62, 0x6B)
SAND = RGBColor(0xC9, 0x7B, 0x3F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def textbox(slide, x, y, w, h, lines, size=18, colour=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    frame = box.text_frame
    frame.word_wrap = True
    for i, line in enumerate(lines):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = colour
        run.font.bold = bold
        run.font.name = "Segoe UI"
        if align is not None:
            para.alignment = align
    return box


def rule(slide, x, y, w, colour=SAND):
    """The one accent, the same 26x3 tick the product's own pages use."""
    from pptx.enum.shapes import MSO_SHAPE

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(Inches(0.04)))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colour
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


# ---------------------------------------------------------------- slide 1: cover
#
# Not part of the template block. It is here so the deck looks like a deck in
# the thumbnail rail rather than like two loose slides, and so the screenshot
# shows a template block sitting INSIDE something, which is how anybody's real
# deck is shaped.
cover = prs.slides.add_slide(blank)
rule(cover, Inches(1.0), Inches(2.6), Inches(0.36))
textbox(cover, Inches(1.0), Inches(2.8), Inches(10), Inches(1.2), ["Quarterly business review"], size=40, bold=True)
textbox(
    cover,
    Inches(1.0),
    Inches(3.9),
    Inches(10),
    Inches(0.8),
    ["FY26 Q1 · one pack per account"],
    size=20,
    colour=GREY,
)

# ------------------------------------------------- slides 2 and 3: the template
#
# Two slides, because one is not enough to show what the product is for: the
# point is that a BLOCK repeats, and a block of one looks like a mail merge for
# labels. Every field name below is a column header in `rows.txt`.
one = prs.slides.add_slide(blank)
rule(one, Inches(0.9), Inches(0.75), Inches(0.36))
textbox(one, Inches(0.9), Inches(0.95), Inches(7.4), Inches(1.0), ["{{Account}} — {{Region}}"], size=34, bold=True)
textbox(
    one,
    Inches(0.9),
    Inches(2.0),
    Inches(5.2),
    Inches(1.6),
    [
        "Renewal {{Renewal|date:d MMM yyyy}}",
        "Annual value {{Revenue|number:0}} EUR",
    ],
    size=18,
    colour=GREY,
)

chart_data = CategoryChartData()
chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
# Plausible rather than round: 18 and 42 is what the kit uses, and it reads as
# test data at a glance.
chart_data.add_series("Revenue, EUR thousands", (284, 301, 297, 368))
chart = one.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Emu(Inches(0.9)),
    Emu(Inches(3.7)),
    Emu(Inches(7.4)),
    Emu(Inches(3.2)),
    chart_data,
).chart
chart.has_legend = False
chart.has_title = True
# The title merges too, which is the feature the description claims and the one
# a static chart would quietly fail to show.
chart.chart_title.text_frame.text = "Quarterly revenue — {{Account}}"
for para in chart.chart_title.text_frame.paragraphs:
    for run in para.runs:
        run.font.size = Pt(14)
        run.font.bold = False
        run.font.color.rgb = GREY

textbox(
    one,
    Inches(9.0),
    Inches(2.0),
    Inches(3.4),
    Inches(4.0),
    [
        "Where it stands",
        "",
        "Renewal conversation opens this quarter. Usage is steady and the "
        "team has asked about the reporting add-on.",
    ],
    size=14,
    colour=GREY,
)

two = prs.slides.add_slide(blank)
rule(two, Inches(0.9), Inches(0.75), Inches(0.36))
textbox(two, Inches(0.9), Inches(0.95), Inches(9), Inches(1.0), ["Next steps for {{Account}}"], size=34, bold=True)
textbox(
    two,
    Inches(0.9),
    Inches(2.2),
    Inches(11.5),
    Inches(3.0),
    [
        "1.  Confirm the renewal date with {{Region}} and put it in the forecast.",
        "",
        "2.  Walk through the reporting add-on before {{Renewal|date:d MMM}}.",
        "",
        "3.  Agree the FY27 number against this year's {{Revenue|number:0}} EUR.",
    ],
    size=20,
)
# Speaker notes merge as well, and nothing else in the shots would show it.
two.notes_slide.notes_text_frame.text = (
    "Owner to send the {{Account}} pack the day after the review. "
    "Renewal is {{Renewal|date:d MMM yyyy}}."
)

prs.save(OUT)

ROWS = os.path.join(HERE, "rows.txt")
with open(ROWS, "w", encoding="utf-8", newline="\n") as f:
    f.write("Account\tRegion\tRevenue\tRenewal\n")
    f.write("Nordwind Retail\tNordics\t1250000\t2026-03-01\n")
    f.write("Brightline Group\tBenelux\t880000\t2026-04-15\n")
    f.write("Alpenhof AG\tDACH\t1640000\t2026-05-30\n")

print(f"demo deck: {os.path.relpath(OUT, os.path.dirname(os.path.dirname(HERE)))}")
print(f"rows:      {os.path.relpath(ROWS, os.path.dirname(os.path.dirname(HERE)))}")
