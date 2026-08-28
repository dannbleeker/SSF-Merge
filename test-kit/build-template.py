"""
Build the SSF Merge test template.

    pip install python-pptx && python3 test-kit/build-template.py

Run rarely, by hand. Not wired into npm or CI: it needs Python and a library
this project does not otherwise use, and the deck it writes is committed beside
it. What CI does run is `test/test-kit.test.ts`, which merges that committed
deck — so the kit cannot rot silently even though nothing rebuilds it.

Written with python-pptx rather than by hand, deliberately: the whole point of a
real-host round is to test the engine against parts somebody ELSE authored. A
chart this repo's own fixture builder wrote would test my writer against my
reader and prove nothing about PowerPoint.

The SmartArt is the one thing left for the user to add in PowerPoint, because
nothing here can author one — and a PowerPoint-authored diagram is the stronger
test anyway.
"""

import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "SSF-Merge-test-template.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

INK = RGBColor(0x11, 0x2B, 0x4A)
GREY = RGBColor(0x55, 0x5F, 0x6D)
ORANGE = RGBColor(0xE4, 0x7A, 0x1C)


def textbox(slide, x, y, w, h, runs, size=18, colour=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    frame = box.text_frame
    frame.word_wrap = True
    for i, line in enumerate(runs):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = colour
        run.font.bold = bold
        run.font.name = "Aptos"
        if align is not None:
            para.alignment = align
    return box


# ---------------------------------------------------------------- slide 1 ----
# Instructions, and NOT part of the block. Everything the round needs is on the
# deck itself, so the file is still usable when this chat is closed.
s1 = prs.slides.add_slide(blank)
textbox(s1, Inches(0.8), Inches(0.5), Inches(11.8), Inches(0.9), ["SSF Merge — test template"], size=34, bold=True)
textbox(
    s1,
    Inches(0.8),
    Inches(1.5),
    Inches(11.8),
    Inches(5.4),
    [
        "The template block is slides 2 and 3. This slide is not part of it.",
        "",
        "1.  Add the SmartArt yourself, on slide 3, where the grey box says so:",
        "     Insert ▸ SmartArt ▸ Process ▸ Basic Process, then type into the three boxes:",
        "     {{Name}}   ·   {{Region}}   ·   Renewal {{Renewal|date:d MMM}}",
        "     Delete the grey box afterwards. (Nothing outside PowerPoint can author SmartArt,",
        "     and a diagram PowerPoint wrote is the stronger test anyway.)",
        "",
        "2.  Open SSF Merge. Step 1: block from 2 to 3. Step 2: paste the rows from data.txt,",
        "     then choose ada.png, grace.png and alan.png when the picture picker appears.",
        "",
        "3.  Step 3 should list every field below. Step 5 merges: 3 rows × 2 slides = 6 slides.",
        "",
        "Fields in this template: Name, Region, Revenue, Renewal, Photo — and Nickname,",
        "which the data deliberately does NOT have. It should stay visible as {{Nickname}}.",
    ],
    size=15,
    colour=GREY,
)

# ---------------------------------------------------------------- slide 2 ----
# Text, a formatted number, a formatted date, a picture frame and a real chart.
s2 = prs.slides.add_slide(blank)
textbox(s2, Inches(0.7), Inches(0.45), Inches(8.0), Inches(0.9), ["{{Name}} — {{Region}}"], size=32, bold=True)
textbox(
    s2,
    Inches(0.7),
    Inches(1.4),
    Inches(6.2),
    Inches(1.2),
    ["Renewal {{Renewal|date:d MMM yyyy}}", "{{Revenue|number:0}} EUR"],
    size=18,
    colour=GREY,
)

# The picture frame. A rectangle whose ONLY text is the image field: the merge
# takes the placeholder away and fills the shape. Deliberately WIDE, so the
# tall photo has to be cropped top and bottom and the wide one barely at all.
frame = s2.shapes.add_textbox(Inches(9.4), Inches(0.7), Inches(3.2), Inches(2.4))
frame.text_frame.word_wrap = True
run = frame.text_frame.paragraphs[0].add_run()
run.text = "{{Photo|image}}"
run.font.size = Pt(12)
run.font.color.rgb = GREY
frame.line.color.rgb = ORANGE
frame.line.width = Pt(1.5)

# A real chart, with placeholders in its title and in its category labels.
# python-pptx writes the cached strings AND the embedded workbook, which is
# exactly the pair the merge has to keep in step.
data = CategoryChartData()
data.categories = ["{{Region}}", "Everyone else"]
data.add_series("Revenue", (18.0, 42.0))
chart_frame = s2.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.7), Inches(2.9), Inches(8.2), Inches(4.1), data
)
chart = chart_frame.chart
chart.has_title = True
chart.chart_title.text_frame.text = "Quarterly revenue — {{Region}}"

s2.notes_slide.notes_text_frame.text = "Call {{Name}} before {{Renewal|date:d MMM}}."

# ---------------------------------------------------------------- slide 3 ----
s3 = prs.slides.add_slide(blank)
textbox(s3, Inches(0.7), Inches(0.45), Inches(11.9), Inches(0.9), ["Next steps for {{Name}}"], size=32, bold=True)
textbox(
    s3,
    Inches(0.7),
    Inches(1.4),
    Inches(11.9),
    Inches(0.6),
    ["Account owner: {{Nickname}}   ← no such column; this should stay as written"],
    size=16,
    colour=GREY,
)

marker = s3.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(11.9), Inches(4.2))
marker.text_frame.word_wrap = True
for i, line in enumerate(
    [
        "Put the SmartArt here.",
        "",
        "Insert ▸ SmartArt ▸ Process ▸ Basic Process, then type into the three boxes:",
        "{{Name}}     {{Region}}     Renewal {{Renewal|date:d MMM}}",
        "",
        "Then delete this box.",
    ]
):
    para = marker.text_frame.paragraphs[0] if i == 0 else marker.text_frame.add_paragraph()
    run = para.add_run()
    run.text = line
    run.font.size = Pt(16)
    run.font.color.rgb = GREY
marker.line.color.rgb = GREY
marker.line.width = Pt(1)

prs.save(OUT)
print("wrote", OUT)
